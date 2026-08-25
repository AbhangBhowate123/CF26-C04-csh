from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation():
    prs = Presentation()

    # Define some standard layouts
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    # Slide 1: Title Slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "CODE FORGE\nAegis Mission Control"
    
    subtitle.text = (
        "Team Name: [CSH]\n"
        "Problem Code: C04\n"
        "Problem Statement Title: Spatial Cyber Threat Reconstruction Engine\n"
        "Team Members: Ashutosh Singrole\n Abhang Bhowate\n"
        "College / Institution: IIIT Nagpur"
    )

    # Function to add a content slide
    def add_content_slide(title_text, content_list):
        slide = prs.slides.add_slide(bullet_slide_layout)
        title = slide.shapes.title
        title.text = title_text
        
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        
        for p_idx, text in enumerate(content_list):
            p = tf.add_paragraph()
            p.text = text
            if ":" in text and "http" not in text:
                parts = text.split(":", 1)
                p.text = ""
                run1 = p.add_run()
                run1.text = parts[0] + ":"
                run1.font.bold = True
                run2 = p.add_run()
                run2.text = parts[1]
            p.level = 0
            if text.strip().startswith("- "):
                p.level = 1

    # Slide 2: The Problem
    add_content_slide("The Problem", [
        "Problem in your own words: Modern buildings contain hundreds of interconnected devices. Security events are often observed independently, making it difficult to detect coordinated attacks across physical locations and network segments.",
        "Real-world context: A sequence of individually harmless events (logins, badge swipes, sensor pings) may represent an attack when viewed together.",
        "Who is affected?: Security analysts, building administrators, enterprise networks.",
        "Why existing approaches are insufficient: They often rely on single opaque alerts without context on how incidents evolve across physical and network relationships.",
        "One clear problem statement: Fragmented device telemetry in smart buildings needs to be correlated into explainable, reconstructed attack paths."
    ])

    # Slide 3: Problem Decomposition
    add_content_slide("Problem Decomposition", [
        "Core technical problem: Correlating events across multiple domains (physical and network) in real-time amidst noisy background telemetry.",
        "Assumptions: Telemetry is available but may be incomplete, duplicated, delayed, or out of order.",
        "Known constraints: Software-only simulation (no real network scanning/intrusion).",
        "Important edge cases: Differentiating legitimate administrative actions from lateral movement; handling delayed or duplicated logs.",
        "What makes the problem difficult?: High volume of background noise masking a subtle, multi-step attack chain across different network segments and floors."
    ])

    # Slide 4: Research & Prior Art (Optional)
    add_content_slide("Research & Prior Art (Optional)", [
        "Existing approaches: Traditional SIEMs, rule-based alerts, isolated physical security systems.",
        "Relevant papers / systems / methods: Graph-based threat detection, kill-chain analysis.",
        "What they do well: Detect known single-point anomalies.",
        "Where they fail for this problem: Lack cross-domain (physical + cyber) context; produce unexplained black-box probability scores.",
        "Your identified gap: A need for transparent, explainable scoring of candidate lateral-movement chains across physical building topologies."
    ])

    # Slide 5: Proposed Approach
    add_content_slide("Proposed Approach", [
        "Solution concept: 'Aegis Mission Control' - A simulated multi-floor building security platform.",
        "Key idea: Correlate fragmented telemetry into candidate attack paths using a sliding time-window graph analysis.",
        "Major components: Telemetry Generator, Flask Backend (Graph + Correlation Engine), Next.js Frontend (Dashboard).",
        "What is novel/different: Explainable, transparent scoring based on human-readable factors (cross-floor movement, privilege escalation) rather than a black box.",
        "Core workflow: Generate synthetic building telemetry -> construct device relationship graph -> run correlation pass -> render reconstructed path and score on live dashboard."
    ])

    # Slide 6: System Architecture
    add_content_slide("System Architecture", [
        "Components: Telemetry Generator (generate_telemetry.py), Backend API (Flask), Frontend (Next.js/React).",
        "Data flow: Telemetry generation -> JSON -> Flask backend constructs graph and scores paths -> REST endpoint /api/attack-paths -> Next.js frontend rendering.",
        "APIs/services: REST API providing scored candidate attack chains and contributing reasons.",
        "Storage: In-memory/JSON (telemetry.json).",
        "Models/algorithms: Sliding time-window correlation over a NetworkX device graph.",
        "External dependencies: Python, Flask, NetworkX, Node.js, Next.js, Tailwind CSS.",
        "Security boundaries where relevant: Logically separates physical zones (floors, VLANs) in the simulation."
    ])

    # Slide 7: Innovation & Impact
    add_content_slide("Innovation & Impact", [
        "Innovation:",
        "  - What is genuinely different?: Transparent, multi-domain (physical + cyber) threat reconstruction.",
        "  - Why is it better?: Provides analysts with an explainable picture (graph, reasons, timeline) instead of a single alert.",
        "  - What technical insight did you introduce?: Modeling physical floor adjacency and network topology together as a unified graph for lateral movement detection.",
        "Impact:",
        "  - Who would use it?: Security Operations Center (SOC) analysts, facility managers.",
        "  - Where could it be deployed?: Smart buildings, corporate campuses, critical infrastructure.",
        "  - What measurable benefit could it create?: Reduced time-to-understand (MTTR) for complex incidents, lower false positive fatigue.",
        "  - What would be required to move beyond prototype?: Integration with real building sensors/logs and ML-assisted scoring."
    ])

    # Slide 8: References (Optional)
    add_content_slide("References (Optional)", [
        "Tools & Frameworks:",
        "  - Python (Flask, NetworkX) for backend correlation",
        "  - Node.js (Next.js, React, Tailwind CSS) for frontend",
        "  - Google Stitch for UI/UX design generation",
        "  - Google Antigravity (Agentic IDE)",
        "  - Claude (Anthropic) for architecture planning and documentation"
    ])

    prs.save('Project_Presentation.pptx')
    print("Presentation saved as Project_Presentation.pptx")

if __name__ == '__main__':
    create_presentation()
